import os
import logging
import magic
from pathlib import Path
from typing import Optional, List, Dict, Any
from fastapi import UploadFile, HTTPException, status
import aiofiles
import aiofiles.os
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation
import csv
from PIL import Image
import pytesseract
from io import BytesIO
import json
import xml.etree.ElementTree as ET

logger = logging.getLogger(__name__)

class FileProcessor:
    SUPPORTED_MIME_TYPES = {
        'application/pdf': 'pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
        'text/csv': 'csv',
        'text/plain': 'txt',
        'application/json': 'json',
        'application/vnd.ms-powerpoint': 'ppt',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
        'application/vnd.ms-excel': 'xls',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
        'image/jpeg': 'jpg',
        'image/png': 'png',
        'application/xml': 'xml',
        'text/xml': 'xml'
    }

    @staticmethod
    async def save_upload_file(upload_dir: str, file: UploadFile) -> str:
        """Guarda archivo subido y valida su tipo"""
        try:
            Path(upload_dir).mkdir(parents=True, exist_ok=True)
            file_path = os.path.join(upload_dir, file.filename)
            
            # Validar tipo de archivo
            file_type = await FileProcessor._validate_file_type(file)
            
            async with aiofiles.open(file_path, 'wb') as buffer:
                await buffer.write(await file.read())
            
            return file_path
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Error saving file: {str(e)}")
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Error saving file: {str(e)}"
            )

    @staticmethod
    async def _validate_file_type(file: UploadFile) -> str:
        """Valida el tipo de archivo usando magic"""
        # Leer los primeros bytes para la detección
        file_content = await file.read(2048)
        await file.seek(0)  # Rewind after reading
        
        mime = magic.from_buffer(file_content, mime=True)
        file_extension = FileProcessor.SUPPORTED_MIME_TYPES.get(mime)
        
        if not file_extension:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Unsupported file type: {mime}"
            )
        
        # Validar extensión vs contenido
        if Path(file.filename).suffix.lower()[1:] != file_extension:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="File extension doesn't match content"
            )
        
        return file_extension

    @staticmethod
    async def extract_text(file_path: str) -> str:
        """Extrae texto de varios formatos de archivo"""
        try:
            file_extension = Path(file_path).suffix.lower()[1:]
            
            handlers = {
                'pdf': FileProcessor._extract_from_pdf,
                'docx': FileProcessor._extract_from_docx,
                'csv': FileProcessor._extract_from_csv,
                'txt': FileProcessor._extract_from_txt,
                'json': FileProcessor._extract_from_json,
                'pptx': FileProcessor._extract_from_pptx,
                'xlsx': FileProcessor._extract_from_xlsx,
                'jpg': FileProcessor._extract_from_image,
                'png': FileProcessor._extract_from_image,
                'xml': FileProcessor._extract_from_xml
            }
            
            if file_extension not in handlers:
                raise ValueError(f"Unsupported file extension: {file_extension}")
            
            return await handlers[file_extension](file_path)
            
        except Exception as e:
            logger.error(f"Error extracting text: {str(e)}")
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Error extracting text: {str(e)}"
            )

    @staticmethod
    async def _extract_from_pdf(file_path: str) -> str:
        """Extrae texto de PDFs"""
        text = ""
        async with aiofiles.open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(await file.read())
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text

    @staticmethod
    async def _extract_from_docx(file_path: str) -> str:
        """Extrae texto de documentos Word"""
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    @staticmethod
    async def _extract_from_csv(file_path: str) -> str:
        """Extrae texto de CSV"""
        async with aiofiles.open(file_path, 'r') as file:
            content = await file.read()
            df = pd.read_csv(BytesIO(content.encode()))
            return df.to_string(index=False)

    @staticmethod
    async def _extract_from_txt(file_path: str) -> str:
        """Lee archivos de texto plano"""
        async with aiofiles.open(file_path, 'r') as file:
            return await file.read()

    @staticmethod
    async def _extract_from_json(file_path: str) -> str:
        """Extrae texto de JSON"""
        async with aiofiles.open(file_path, 'r') as file:
            data = json.loads(await file.read())
            return str(data)

    @staticmethod
    async def _extract_from_pptx(file_path: str) -> str:
        """Extrae texto de PowerPoint"""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    @staticmethod
    async def _extract_from_xlsx(file_path: str) -> str:
        """Extrae texto de Excel"""
        df = pd.read_excel(file_path, engine='openpyxl')
        return df.to_string(index=False)

    @staticmethod
    async def _extract_from_image(file_path: str) -> str:
        """Extrae texto de imágenes usando OCR"""
        async with aiofiles.open(file_path, 'rb') as file:
            image = Image.open(BytesIO(await file.read()))
            return pytesseract.image_to_string(image)

    @staticmethod
    async def _extract_from_xml(file_path: str) -> str:
        """Extrae texto de XML"""
        async with aiofiles.open(file_path, 'r') as file:
            tree = ET.parse(BytesIO((await file.read()).encode()))
            root = tree.getroot()
            return ET.tostring(root, encoding='unicode', method='text')

    @staticmethod
    async def clean_up(file_path: str) -> None:
        """Elimina archivos temporales"""
        try:
            if await aiofiles.os.path.exists(file_path):
                await aiofiles.os.remove(file_path)
        except Exception as e:
            logger.warning(f"Could not delete file {file_path}: {str(e)}")